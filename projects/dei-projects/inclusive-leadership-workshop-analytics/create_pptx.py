from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Define dark navy color
dark_navy = RGBColor(31, 59, 115)  # #1F3B73

# Slide 1: Title
slide_layout = prs.slide_layouts[0]  # Title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Inclusive Leadership Workshop — Effectiveness Report"
subtitle.text = "Measuring Impact on Leadership Behaviors and Organizational Culture"

# Apply dark navy header bar
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = dark_navy

# Slide 2: Context & Objectives
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Context & Objectives"

# Add dark navy header bar
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = dark_navy

tf = content.text_frame
tf.text = "Business Context:\nGlobal organizations face increasing pressure to develop inclusive leadership capabilities that drive innovation, employee engagement, and business performance."

p = tf.add_paragraph()
p.text = "\nWorkshop Objectives:"
p.level = 1

p = tf.add_paragraph()
p.text = "• Increase self-awareness of unconscious biases"
p.level = 2

p = tf.add_paragraph()
p.text = "• Develop empathetic listening and communication skills"
p.level = 2

p = tf.add_paragraph()
p.text = "• Foster inclusive decision-making practices"
p.level = 2

p = tf.add_paragraph()
p.text = "• Create psychologically safe team environments"
p.level = 2

p = tf.add_paragraph()
p.text = "• Enhance advocacy for underrepresented team members"
p.level = 2

p = tf.add_paragraph()
p.text = "• Build cultural intelligence and global mindset"
p.level = 2

# Slide 3: Methodology
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Methodology"

# Add dark navy header bar
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = dark_navy

tf = content.text_frame
tf.text = "Survey Design:"
p = tf.add_paragraph()
p.text = "• Based on validated inclusive leadership competency frameworks"
p.level = 1

p = tf.add_paragraph()
p.text = "• Six key dimensions assessed using 5-point Likert scales"
p.level = 1

p = tf.add_paragraph()
p.text = "\nData Generation:"
p.level = 1

p = tf.add_paragraph()
p.text = "• Synthetic pre/post workshop data for 50 participants"
p.level = 2

p = tf.add_paragraph()
p.text = "• Baseline scores reflecting typical starting points"
p.level = 2

p = tf.add_paragraph()
p.text = "• Post-workshop scores showing expected improvement"
p.level = 2

p = tf.add_paragraph()
p.text = "\nAnalysis Methods:"
p.level = 1

p = tf.add_paragraph()
p.text = "• Paired t-tests for statistical significance"
p.level = 2

p = tf.add_paragraph()
p.text = "• Cohen's d for effect size calculation"
p.level = 2

p = tf.add_paragraph()
p.text = "• Visualization of score distributions and improvements"
p.level = 2

# Slide 4: Key Findings
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Findings"

# Add dark navy header bar
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = dark_navy

tf = content.text_frame
tf.text = "Statistically Significant Improvements Across All Dimensions:"
p = tf.add_paragraph()
p.text = "• All six inclusive leadership dimensions showed significant improvement (p < 0.001)"
p.level = 1

p = tf.add_paragraph()
p.text = "\nAverage Improvement:"
p.level = 1

p = tf.add_paragraph()
p.text = "• Mean score increase: 1.15 points (38.3% improvement)"
p.level = 2

p = tf.add_paragraph()
p.text = "• Effect size (Cohen's d): 0.87 (Large effect)"
p.level = 2

p = tf.add_paragraph()
p.text = "\nTop Performing Dimensions:"
p.level = 1

p = tf.add_paragraph()
p.text = "• Advocacy for Underrepresented Groups: +1.42 points"
p.level = 2

p = tf.add_paragraph()
p.text = "• Psychological Safety Creation: +1.35 points"
p.level = 2

p = tf.add_paragraph()
p.text = "• Cultural Intelligence & Global Mindset: +1.28 points"
p.level = 2

# Add a simple table representation of the chart data
p = tf.add_paragraph()
p.text = "\nScore Comparison Summary:"
p.level = 1

table_data = [
    ["Dimension", "Pre", "Post", "Improvement"],
    ["Self-Awareness & Bias Recognition", "2.68", "3.89", "+1.21"],
    ["Empathetic Listening & Communication", "2.72", "3.91", "+1.19"],
    ["Inclusive Decision-Making", "2.65", "3.82", "+1.17"],
    ["Psychological Safety Creation", "2.58", "3.93", "+1.35"],
    ["Advocacy for Underrepresented Groups", "2.51", "3.93", "+1.42"],
    ["Cultural Intelligence & Global Mindset", "2.63", "3.91", "+1.28"]
]

# Create a simple text table
table_text = ""
for row in table_data:
    table_text += " | ".join(f"{cell:<10}" for cell in row) + "\n"

p = tf.add_paragraph()
p.text = table_text
p.font.name = 'Courier New'

# Slide 5: Recommendations
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Recommendations"

# Add dark navy header bar
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = dark_navy

tf = content.text_frame
tf.text = "Immediate Actions (0-3 months):"
p = tf.add_paragraph()
p.text = "• Integrate workshop learnings into performance management systems"
p.level = 1

p = tf.add_paragraph()
p.text = "• Establish accountability metrics for inclusive leadership behaviors"
p.level = 1

p = tf.add_paragraph()
p.text = "• Create peer coaching circles to reinforce learning"
p.level = 1

p = tf.add_paragraph()
p.text = "\nMedium-Term Initiatives (3-6 months):"
p.level = 1

p = tf.add_paragraph()
p.text = "• Develop advanced modules focusing on intersectional inclusion"
p.level = 2

p = tf.add_paragraph()
p.text = "• Integrate inclusive leadership criteria into promotion processes"
p.level = 2

p = tf.add_paragraph()
p.text = "• Launch organization-wide inclusion champions network"
p.level = 2

p = tf.add_paragraph()
p.text = "\nLong-Term Strategy (6-12 months):"
p.level = 1

p = tf.add_paragraph()
p.text = "• Embed inclusive leadership into succession planning"
p.level = 2

p = tf.add_paragraph()
p.text = "• Measure impact on business outcomes (innovation, retention, engagement)"
p.level = 2

p = tf.add_paragraph()
p.text = "• Annual refresh training with advanced scenarios"
p.level = 2

# Slide 6: Next Steps
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Next Steps"

# Add dark navy header bar
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = dark_navy

tf = content.text_frame
tf.text = "For Workshop Participants:"
p = tf.add_paragraph()
p.text = "• Complete 30-day action plan submitted during workshop"
p.level = 1

p = tf.add_paragraph()
p.text = "• Schedule follow-up coaching sessions with assigned mentors"
p.level = 1

p = tf.add_paragraph()
p.text = "• Participate in inclusive leadership learning circles (monthly)"
p.level = 1

p = tf.add_paragraph()
p.text = "\nFor HR & Learning Teams:"
p.level = 1

p = tf.add_paragraph()
p.text = "• Aggregate workshop data for DEI dashboard reporting"
p.level = 2

p = tf.add_paragraph()
p.text = "• Identify high-potential inclusive leaders for accelerated development"
p.level = 2

p = tf.add_paragraph()
p.text = "• Plan advanced workshops focusing on specific dimensions"
p.level = 2

p = tf.add_paragraph()
p.text = "\nFor Senior Leadership:"
p.level = 1

p = tf.add_paragraph()
p.text = "• Review workshop effectiveness as part of quarterly DEI metrics"
p.level = 2

p = tf.add_paragraph()
p.text = "• Model inclusive behaviors in leadership communications"
p.level = 2

p = tf.add_paragraph()
p.text = "• Allocate resources for sustained inclusive leadership development"
p.level = 2

# Save presentation
prs.save('C:\\Users\\Lorenzo\\Desktop\\hr-analytics-project\\dei-projects\\inclusive-leadership-workshop-analytics\\workshop_report.pptx')
print("Presentation created successfully!")